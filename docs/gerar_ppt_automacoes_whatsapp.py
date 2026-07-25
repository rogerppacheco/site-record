"""
Gera apresentação PPT das automações WhatsApp — Record PAP.
Identidade visual alinhada ao Design System Record PAP v17.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import nsmap
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Emu, Inches, Pt

# --- Cores Record PAP v17 ---
ACCENT = RGBColor(0x0E, 0xA5, 0xE9)
ACCENT_HOVER = RGBColor(0x02, 0x84, 0xC7)
ACCENT_LIGHT = RGBColor(0x38, 0xBD, 0xF8)
SECONDARY = RGBColor(0x06, 0xB6, 0xD4)
SUCCESS = RGBColor(0x05, 0x96, 0x69)
DANGER = RGBColor(0xDC, 0x26, 0x26)
WARNING = RGBColor(0xD9, 0x77, 0x06)
BG = RGBColor(0xF7, 0xF9, 0xFC)
BG_SECONDARY = RGBColor(0xE8, 0xED, 0xF5)
SURFACE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT = RGBColor(0x1A, 0x23, 0x32)
TEXT_MUTED = RGBColor(0x5B, 0x6B, 0x7D)
FOOTER = RGBColor(0x1E, 0x29, 0x3B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TIP_BG = RGBColor(0xE0, 0xF2, 0xFE)
WARN_BG = RGBColor(0xFE, 0xF3, 0xC7)
OK_BG = RGBColor(0xD1, 0xFA, 0xE5)
DANGER_BG = RGBColor(0xFE, 0xE2, 0xE2)

ROOT = Path(__file__).resolve().parents[1]
LOGO = ROOT / "static" / "logo.png"
OUT = ROOT / "docs" / "Automacoes_WhatsApp_Record_PAP.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _set_run_font(run, size: int, bold: bool = False, color: RGBColor = TEXT, name: str = "Calibri") -> None:
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name


def _fill_solid(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _add_rect(slide, left, top, width, height, color: RGBColor):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    _fill_solid(shape, color)
    return shape


def _add_round_rect(slide, left, top, width, height, color: RGBColor):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    _fill_solid(shape, color)
    # cantos um pouco mais suaves
    try:
        shape.adjustments[0] = 0.08
    except Exception:
        pass
    return shape


def _textbox(slide, left, top, width, height, text: str, size: int = 14, bold: bool = False,
             color: RGBColor = TEXT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf._txBody.bodyPr.set("anchor", "t" if anchor == MSO_ANCHOR.TOP else "ctr")
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_run_font(run, size, bold, color)
    return box


def _add_paragraph(tf, text: str, size: int = 13, bold: bool = False, color: RGBColor = TEXT,
                   space_before: int = 4, space_after: int = 2, align=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    _set_run_font(run, size, bold, color)
    return p


def _base_slide(prs: Presentation, with_footer: bool = True):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, BG)
    # barra superior
    _add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), ACCENT)
    if with_footer:
        _add_rect(slide, 0, Inches(7.15), SLIDE_W, Inches(0.35), FOOTER)
        _textbox(
            slide, Inches(0.4), Inches(7.18), Inches(8), Inches(0.28),
            "Record PAP  ·  Automações WhatsApp  ·  Uso interno",
            size=10, color=RGBColor(0x94, 0xA3, 0xB8),
        )
        _textbox(
            slide, Inches(10.5), Inches(7.18), Inches(2.5), Inches(0.28),
            "Confidencial",
            size=10, color=RGBColor(0x94, 0xA3, 0xB8), align=PP_ALIGN.RIGHT,
        )
    return slide


def _header(slide, title: str, subtitle: str = "", logo: bool = True):
    if logo and LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(0.4), Inches(0.22), height=Inches(0.55))
    _textbox(slide, Inches(2.2) if logo else Inches(0.4), Inches(0.22), Inches(10), Inches(0.4),
             title, size=24, bold=True, color=TEXT)
    if subtitle:
        _textbox(slide, Inches(2.2) if logo else Inches(0.4), Inches(0.62), Inches(10.5), Inches(0.3),
                 subtitle, size=12, color=TEXT_MUTED)


def _chip(slide, left, top, width, height, text: str, bg: RGBColor, fg: RGBColor):
    _add_round_rect(slide, left, top, width, height, bg)
    _textbox(slide, left, top + Inches(0.05), width, height - Inches(0.05),
             text, size=11, bold=True, color=fg, align=PP_ALIGN.CENTER)


def _card(slide, left, top, width, height, title: str, lines: list[str],
          title_color: RGBColor = ACCENT_HOVER, accent_bar: RGBColor | None = ACCENT):
    _add_round_rect(slide, left, top, width, height, SURFACE)
    if accent_bar:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.08), height)
        _fill_solid(bar, accent_bar)
    box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.12), width - Inches(0.35), height - Inches(0.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    _set_run_font(run, 13, True, title_color)
    for line in lines:
        _add_paragraph(tf, line, size=12, color=TEXT, space_before=3, space_after=1)


def _phone_mock(slide, left, top, width, height, title: str, bubbles: list[tuple[str, str]]):
    """bolhas: ('bot'|'user', texto)"""
    _add_round_rect(slide, left, top, width, height, FOOTER)
    # header do "celular"
    _add_rect(slide, left, top, width, Inches(0.42), ACCENT_HOVER)
    _textbox(slide, left, top + Inches(0.08), width, Inches(0.3),
             title, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    y = top + Inches(0.55)
    for who, msg in bubbles:
        is_bot = who == "bot"
        bubble_w = width - Inches(0.45)
        # estima altura pelo tamanho do texto
        lines_est = max(1, (len(msg) // 38) + 1)
        bh = Inches(0.22 + 0.18 * lines_est)
        if y + bh > top + height - Inches(0.15):
            break
        bx = left + Inches(0.15) if is_bot else left + Inches(0.3)
        bg = RGBColor(0x33, 0x41, 0x55) if is_bot else ACCENT
        _add_round_rect(slide, bx, y, bubble_w, bh, bg)
        _textbox(slide, bx + Inches(0.1), y + Inches(0.04), bubble_w - Inches(0.15), bh - Inches(0.05),
                 msg, size=10, color=WHITE)
        y += bh + Inches(0.08)


# ---------------------------------------------------------------------------
# Conteúdo
# ---------------------------------------------------------------------------

AUTOMACOES = [
    {
        "cmd": "DFV",
        "nome": "Consultar fachadas por CEP",
        "tag": "Power BI ao vivo",
        "auth": "Todos os vendedores cadastrados",
        "passos": [
            "Digite DFV no WhatsApp.",
            "Informe o CEP (8 dígitos; hífen opcional).",
            "Aguarde a consulta ao Power BI ao vivo.",
            "Receba logradouro, bairro, cidade/UF, CDO(s) e números/complementos.",
        ],
        "telas": [
            ("bot", "Digite o CEP para consultar fachadas (DFV)."),
            ("user", "29101-400"),
            ("bot", "🏢 DFV (Power BI ao vivo)\nRua X · Vitória/ES\nNúmeros: 10, 12, 14A…"),
        ],
        "dicas": [
            "Prefira números viáveis quando a lista for longa.",
            "Cobertura: ES, MG e RJ.",
            "Digite CANCELAR / SAIR / PARAR para abortar.",
        ],
        "avisos": [
            "Fonte ao vivo do Power BI — pode haver latência.",
            "Fachada antiga (FACHADA) foi desativada; use DFV.",
        ],
    },
    {
        "cmd": "CDOE",
        "nome": "Endereços por código do CDO",
        "tag": "Power BI",
        "auth": "Todos os vendedores cadastrados",
        "passos": [
            "Digite CDOE ou CDOE 28005 (código inline).",
            "Selecione a UF nos botões: MG / ES / RJ.",
            "Se houver várias cidades, escolha na lista interativa.",
            "Escolha o número da rua e veja os números de fachada.",
        ],
        "telas": [
            ("bot", "Informe o código do CDO (ex: 28005)."),
            ("user", "28005"),
            ("bot", "Selecione a UF: [MG] [ES] [RJ]"),
        ],
        "dicas": [
            "Aceita 28005 ou CDOE-28005.",
            "Use a lista de opções do WhatsApp para escolher a rua.",
        ],
        "avisos": [
            "Mesma fonte Power BI do DFV.",
            "Cancele a qualquer momento com CANCELAR.",
        ],
    },
    {
        "cmd": "VIABILIDADE",
        "nome": "Consultar viabilidade por CEP e número",
        "tag": "KMZ / Mapa",
        "auth": "Todos os vendedores cadastrados",
        "passos": [
            "Digite VIABILIDADE.",
            "Informe o CEP.",
            "Informe o número da fachada (Modo Mapa).",
            "O bot geolocaliza e testa se o ponto está na mancha (KMZ).",
        ],
        "telas": [
            ("bot", "Informe o CEP para viabilidade."),
            ("user", "29101400"),
            ("bot", "Agora o número da fachada."),
            ("user", "150"),
        ],
        "dicas": [
            "Se a geo falhar, o sistema tenta a base DFV local.",
            "Aliases: VIABILIDADE / VIABILIDADES.",
        ],
        "avisos": [
            "Resultado sujeito a vistoria técnica local.",
            "Pode retornar: dentro da mancha, fora, não localizado ou só no DFV.",
        ],
    },
    {
        "cmd": "INCLUSAO",
        "nome": "Solicitar inclusão de viabilidade",
        "tag": "Formulário",
        "auth": "Flag autorizar_inclusao_wpp = Sim",
        "passos": [
            "Digite INCLUSAO.",
            "CEP → número (SN vira 0) → complementos → vizinhos (Frente/Dir/Esq).",
            "Confirme coordenadas (Maps) ou cole lat/long manualmente.",
            "Foto (Street View auto ou envie imagem) + comprovantes opcionais.",
            "Observações → confirme com SIM (ética) ou CANCELAR.",
        ],
        "telas": [
            ("bot", "Inclusão de viabilidade. Informe o CEP."),
            ("user", "29101400"),
            ("bot", "Confirme as coordenadas no Maps: SIM ou cole lat,long."),
        ],
        "dicas": [
            "Envie foto real da fachada quando o Street View falhar.",
            "Comprovantes (PDF/imagem) são opcionais — digite PRONTO para seguir.",
        ],
        "avisos": [
            "NÃO invente fachada nem complemento.",
            "Há lock anti-duplicidade; anexos temporários são apagados após o envio.",
            "Requer permissão na Governança.",
        ],
    },
    {
        "cmd": "STATUS",
        "nome": "Consultar status de pedido",
        "tag": "CRM + PAP online",
        "auth": "Todos os vendedores cadastrados",
        "passos": [
            "Digite STATUS (ou SITUACAO).",
            "Escolha: 1 = CPF  ou  2 = O.S.",
            "Informe o dado solicitado.",
            "Receba o status do CRM; se elegível, consulta online no PAP (pode enviar print).",
        ],
        "telas": [
            ("bot", "Consultar status:\n1 — CPF\n2 — O.S."),
            ("user", "1"),
            ("bot", "Informe o CPF (somente números)."),
        ],
        "dicas": [
            "A consulta online pode levar até ~3 minutos.",
            "Use O.S. quando já tiver o número do pedido.",
        ],
        "avisos": [
            "Timeout da consulta online ≈ 180s.",
            "Aguarde a mensagem final antes de digitar outro comando.",
        ],
    },
    {
        "cmd": "FATURA",
        "nome": "Consultar fatura por CPF",
        "tag": "Nio Negociar / API",
        "auth": "Todos os vendedores cadastrados",
        "passos": [
            "Digite FATURA.",
            "Informe o CPF (11 dígitos).",
            "Veja a lista (atrasadas → abertas → outras).",
            "Se houver várias, digite o número da fatura desejada para detalhes + PDF.",
        ],
        "telas": [
            ("bot", "Informe o CPF para consultar faturas."),
            ("user", "12345678901"),
            ("bot", "1) Atrasada R$ 99,90\n2) Aberta R$ 89,90\nDigite o nº…"),
        ],
        "dicas": [
            "O DV do CPF não é validado localmente — a Nio valida.",
            "PDF pode vir da API ou via navegação automatizada.",
        ],
        "avisos": [
            "Evite repetir o comando em sequência (lock anti-duplicidade).",
            "Use Conta quando precisar da 2ª via pelo site Nio.",
        ],
    },
    {
        "cmd": "CONTA",
        "nome": "2ª via de conta por CPF",
        "tag": "Site Nio",
        "auth": "Todos os vendedores cadastrados",
        "passos": [
            "Digite CONTA.",
            "Informe o CPF.",
            "Receba detalhes + PDF da primeira conta encontrada no site.",
        ],
        "telas": [
            ("bot", "Informe o CPF para 2ª via de conta."),
            ("user", "12345678901"),
            ("bot", "📄 Conta encontrada — enviando PDF…"),
        ],
        "dicas": [
            "Fluxo mais direto que FATURA (foco em 2ª via).",
            "Aliases: CONTA / CONTAS.",
        ],
        "avisos": [
            "Depende da disponibilidade do site Nio.",
            "Confirme o CPF digitado antes de enviar.",
        ],
    },
    {
        "cmd": "MATERIAL / APOIA",
        "nome": "Buscar materiais e documentos",
        "tag": "Record Apoia",
        "auth": "Todos os vendedores cadastrados",
        "passos": [
            "Digite MATERIAL ou APOIA (ou Record Apoia).",
            "Informe uma palavra-chave (≥ 2 caracteres).",
            "1 arquivo → envio direto; vários → lista numerada para escolher.",
            "Dica: na etapa inicial, texto livre também busca por tag.",
        ],
        "telas": [
            ("bot", "📁 Record Apoia — digite a palavra-chave."),
            ("user", "plano"),
            ("bot", "1) Tabela Planos\n2) Script Venda\nDigite o número…"),
        ],
        "dicas": [
            "Busca em tags, título, descrição e categoria.",
            "PDF > 5MB sobe para link (R2) em vez do arquivo no chat.",
        ],
        "avisos": [
            "Arquivos inativos não aparecem.",
            "Peça ao admin para reenviar material se o arquivo estiver corrompido.",
        ],
    },
    {
        "cmd": "ANDAMENTO",
        "nome": "Ver agendamentos do dia",
        "tag": "1 passo",
        "auth": "Vendedor: só as próprias O.S. · Diretoria/Admin/BackOffice: todas",
        "passos": [
            "Digite ANDAMENTO.",
            "Receba a lista do dia com cliente, O.S., horário, município e atividade.",
        ],
        "telas": [
            ("user", "ANDAMENTO"),
            ("bot", "📅 Agendamentos de hoje\n• Maria — OS 123 — 14h — Vitória — Instalação"),
        ],
        "dicas": [
            "Ideal para acompanhar a operação do dia em segundos.",
            "Aliases: ANDAMENTO / ANDAMENTOS.",
        ],
        "avisos": [
            "Lista apenas O.S. com horário de execução real preenchido.",
            "Vendedor comum não vê agenda de outros.",
        ],
    },
    {
        "cmd": "CREDITO",
        "nome": "Análise de crédito por CPF/CNPJ",
        "tag": "PAP BackOffice",
        "auth": "Flag autorizar_analise_credito_wpp = Sim",
        "passos": [
            "Digite CREDITO.",
            "Informe CPF (11) ou CNPJ (14).",
            "Se CNPJ, informe também o CPF do representante.",
            "Aguarde o resultado em background (APROVADO / NEGADO + formas de pgto).",
        ],
        "telas": [
            ("bot", "Informe CPF ou CNPJ para análise de crédito."),
            ("user", "12345678901"),
            ("bot", "✅ Análise: APROVADO\nFormas: …"),
        ],
        "dicas": [
            "Histórico fica registrado no sistema.",
            "Use antes de avançar uma venda complexa.",
        ],
        "avisos": [
            "Limite: 1 análise por minuto e 15 por dia (horário de Brasília).",
            "Requer permissão na Governança + pool BO PAP disponível.",
        ],
    },
    {
        "cmd": "PEDIDO",
        "nome": "Consultar pedido / O.S. por CPF no PAP",
        "tag": "Últimos 30 dias + print",
        "auth": "Mesma permissão de Crédito (autorizar_analise_credito_wpp)",
        "passos": [
            "Digite PEDIDO.",
            "Informe CPF ou CNPJ.",
            "Receba texto (status/data/plano/OS) + screenshot da lista no PAP.",
        ],
        "telas": [
            ("bot", "Informe o CPF/CNPJ para consultar pedidos."),
            ("user", "12345678901"),
            ("bot", "📦 Pedidos (30 dias)\nOS 999 — Instalação — Plano X + print"),
        ],
        "dicas": [
            "Janela de consulta: últimos 30 dias.",
            "O print ajuda a validar visualmente no PAP.",
        ],
        "avisos": [
            "Sem resultado: mensagem “Não tem pedido com 30 dias…” + imagem.",
            "Depende do pool de logins BackOffice.",
        ],
    },
    {
        "cmd": "BIO",
        "nome": "Resultado da biometria (Br Pronto)",
        "tag": "GED / Br Pronto",
        "auth": "Flag autorizar_consulta_bio_wpp (Governança)",
        "passos": [
            "Digite BIO ou BIOMETRIA.",
            "Informe o CPF.",
            "Receba se está apto (“Doc. Apto para Venda”) ou os status encontrados + print.",
        ],
        "telas": [
            ("bot", "Informe o CPF para consultar biometria."),
            ("user", "12345678901"),
            ("bot", "✅ Doc. Apto para Venda + print"),
        ],
        "dicas": [
            "Use antes de concluir vendas que exigem biometria.",
            "O print confirma o que aparece no GED.",
        ],
        "avisos": [
            "Logoff obrigatório no GED após a consulta (feito pela automação).",
            "Somente quem tem a flag na Governança.",
        ],
    },
    {
        "cmd": "VENDER",
        "nome": "Realizar venda pelo WhatsApp",
        "tag": "PAP completo 🆕",
        "auth": "autorizar_venda_sem_auditoria + matrícula PAP",
        "passos": [
            "Digite VENDER (ou VENDA) e confirme a matrícula (SIM).",
            "CEP → número → referência (viabilidade no PAP).",
            "CPF → celular → e-mail → crédito → forma pgto → plano → stream/fixo.",
            "Cliente confirma SIM no WhatsApp + biometria (CONSULTAR / BIO OK).",
            "Agendamento (dia/período) → Abrir O.S. → grava no CRM.",
        ],
        "telas": [
            ("bot", "Confirma matrícula PAP xxx? [SIM]"),
            ("user", "SIM"),
            ("bot", "Informe o CEP da instalação…"),
        ],
        "dicas": [
            "ESTENDER: +5 min (máx. 3 vezes). REPETIR: reenvia última pergunta.",
            "CANCELAR / SAIR / PARAR abortam o fluxo com segurança.",
            "Timeout de sessão ≈ 10 minutos de inatividade.",
        ],
        "avisos": [
            "NOVA VENDA é outro fluxo (cadastro CRM), não abre o PAP.",
            "Exige pool BO PAP saudável e matrícula válida.",
            "Confirmação do cliente no WhatsApp é obrigatória.",
        ],
    },
    {
        "cmd": "NOVA VENDA",
        "nome": "Cadastrar venda no CRM",
        "tag": "Via APP ou Sem APP",
        "auth": "Vendedores cadastrados (automática: flag autorizar_venda_automatica)",
        "passos": [
            "Digite NOVA VENDA (ou CADASTRAR VENDA).",
            "Escolha APP ou SEM APP → telefone fixo SIM/NÃO.",
            "Se autorizado: informe se gerou O.S. automática.",
            "CPF/CNPJ → nome (se novo) → tel1 → tel2.",
            "Via APP: vai a observações. Sem APP: endereço, plano, forma pgto, obs.",
            "Confirme e salve no CRM.",
        ],
        "telas": [
            ("bot", "Nova venda no CRM:\n[APP] ou [SEM APP]?"),
            ("user", "SEM APP"),
            ("bot", "Informe o CPF/CNPJ do cliente…"),
        ],
        "dicas": [
            "Use APP quando a venda já foi feita no app Nio.",
            "SEM APP preenche endereço e plano para o time operar.",
        ],
        "avisos": [
            "Não confunda com VENDER (que abre pedido no PAP).",
            "CNPJ pede dados extras (mãe, nascimento, e-mail, representante).",
        ],
    },
    {
        "cmd": "COMISSAO",
        "nome": "Bônus, desconto e adiantamentos",
        "tag": "Somente Diretoria / Admin",
        "auth": "Perfil Diretoria ou Admin",
        "passos": [
            "Digite COMISSAO para ver a ajuda dos comandos.",
            "BONUS login valor AAAA-MM-DD descrição",
            "DESCONTO login valor AAAA-MM-DD descrição",
            "ADIANT_COMISSAO [AAAA-MM] [login] → LISTA / DETALHE N / marcar vendas",
            "ADIANT_SABADO [data] → resumo do sábado + mesma interação.",
        ],
        "telas": [
            ("user", "COMISSAO"),
            ("bot", "Ajuda: BONUS · DESCONTO · ADIANT_COMISSAO · ADIANT_SABADO"),
            ("user", "BONUS joao 50 2026-07-22 Meta batida"),
        ],
        "dicas": [
            "Siga exatamente o formato data AAAA-MM-DD.",
            "Em adiantamentos, use LISTA e DETALHE N para navegar.",
        ],
        "avisos": [
            "Comandos sensíveis — apenas Diretoria/Admin.",
            "Valide login e valores antes de confirmar lançamentos.",
        ],
    },
]


def slide_capa(prs: Presentation) -> None:
    slide = _base_slide(prs, with_footer=False)
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, FOOTER)
    # faixa accent
    _add_rect(slide, 0, 0, Inches(0.25), SLIDE_H, ACCENT)
    _add_rect(slide, 0, Inches(5.9), SLIDE_W, Inches(1.6), ACCENT_HOVER)

    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(0.7), Inches(1.0), height=Inches(0.85))

    _textbox(slide, Inches(0.7), Inches(2.2), Inches(11), Inches(0.7),
             "Automações WhatsApp", size=40, bold=True, color=WHITE)
    _textbox(slide, Inches(0.7), Inches(2.95), Inches(11), Inches(0.45),
             "Guia prático · Passo a passo · Dicas e avisos", size=20, color=ACCENT_LIGHT)
    _textbox(slide, Inches(0.7), Inches(3.6), Inches(10), Inches(0.8),
             "Record PAP — treinamento interno para o time comercial\n"
             "Menu completo: DFV, CDOE, Viabilidade, Inclusão, Status, Fatura, Conta,\n"
             "Material/Apoia, Andamento, Crédito, Pedido, Bio, Vender, Nova Venda e Comissão.",
             size=14, color=RGBColor(0xCB, 0xD5, 0xE1))
    _textbox(slide, Inches(0.7), Inches(6.2), Inches(10), Inches(0.5),
             "Digite MENU no WhatsApp para ver todas as opções",
             size=16, bold=True, color=WHITE)
    _textbox(slide, Inches(0.7), Inches(6.7), Inches(10), Inches(0.35),
             "Identidade visual: Design System Record PAP v17",
             size=11, color=RGBColor(0xE0, 0xF2, 0xFE))


def slide_agenda(prs: Presentation) -> None:
    slide = _base_slide(prs)
    _header(slide, "Agenda do treinamento", "16 automações do menu WhatsApp")

    cols = [
        AUTOMACOES[:8],
        AUTOMACOES[8:],
    ]
    for ci, grupo in enumerate(cols):
        left = Inches(0.4 + ci * 6.4)
        top = Inches(1.15)
        _add_round_rect(slide, left, top, Inches(6.1), Inches(5.7), SURFACE)
        box = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.2), Inches(5.6), Inches(5.3))
        tf = box.text_frame
        tf.word_wrap = True
        first = True
        for i, a in enumerate(grupo, start=1 + ci * 8):
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            p.space_before = Pt(6)
            p.space_after = Pt(2)
            run = p.add_run()
            run.text = f"{i:02d}.  {a['cmd']}"
            _set_run_font(run, 13, True, ACCENT_HOVER)
            p2 = tf.add_paragraph()
            p2.space_before = Pt(0)
            p2.space_after = Pt(4)
            run2 = p2.add_run()
            run2.text = f"     {a['nome']}"
            _set_run_font(run2, 11, False, TEXT_MUTED)


def slide_pre_requisitos(prs: Presentation) -> None:
    slide = _base_slide(prs)
    _header(slide, "Antes de começar", "Pré-requisitos e regras gerais")

    _card(
        slide, Inches(0.4), Inches(1.15), Inches(6.1), Inches(2.6),
        "🔑 Acesso ao menu de vendedor",
        [
            "• Telefone precisa estar cadastrado em Usuario.tel_whatsapp",
            "  (também _2 / _3) e o usuário deve estar ativo.",
            "• Número não cadastrado → atendimento de cliente/IA,",
            "  não o menu comercial.",
            "• Digite MENU, AJUDA, HELP ou OPCOES para reabrir o menu.",
        ],
        accent_bar=ACCENT,
    )
    _card(
        slide, Inches(6.8), Inches(1.15), Inches(6.1), Inches(2.6),
        "🧭 Como navegar nos fluxos",
        [
            "• Sempre um comando por vez (ex.: DFV).",
            "• CANCELAR / SAIR / PARAR abortam a maioria dos fluxos.",
            "• Aguarde a resposta do bot antes de digitar o próximo dado.",
            "• Algumas opções exigem flag na Governança.",
            "• Sessões longas (ex.: VENDER) têm timeout de inatividade.",
        ],
        accent_bar=SECONDARY,
    )
    _card(
        slide, Inches(0.4), Inches(4.0), Inches(6.1), Inches(2.8),
        "🛡️ Permissões especiais (Governança)",
        [
            "• Inclusão → autorizar_inclusao_wpp",
            "• Crédito / Pedido → autorizar_analise_credito_wpp",
            "• Bio → autorizar_consulta_bio_wpp",
            "• Vender → autorizar_venda_sem_auditoria + matrícula PAP",
            "• Comissão → somente Diretoria ou Admin",
        ],
        title_color=WARNING,
        accent_bar=WARNING,
    )
    _card(
        slide, Inches(6.8), Inches(4.0), Inches(6.1), Inches(2.8),
        "💡 Boas práticas",
        [
            "• Confirme CEP e CPF antes de enviar.",
            "• Não invente fachada/complemento na Inclusão.",
            "• Em dúvidas de material, use tags curtas (ex.: plano).",
            "• VENDER ≠ NOVA VENDA (PAP vs cadastro CRM).",
            "• Em erro persistente, avise o BackOffice / TI.",
        ],
        title_color=SUCCESS,
        accent_bar=SUCCESS,
    )


def slide_automacao(prs: Presentation, idx: int, a: dict) -> None:
    slide = _base_slide(prs)
    _header(
        slide,
        f"{idx:02d}. {a['cmd']} — {a['nome']}",
        f"{a['tag']}  ·  Acesso: {a['auth']}",
    )

    # chips
    _chip(slide, Inches(0.4), Inches(1.05), Inches(2.2), Inches(0.32),
          f"Comando: {a['cmd'].split('/')[0].strip()}", TIP_BG, ACCENT_HOVER)
    _chip(slide, Inches(2.75), Inches(1.05), Inches(2.0), Inches(0.32),
          a["tag"][:28], OK_BG, SUCCESS)

    # Passo a passo
    _add_round_rect(slide, Inches(0.4), Inches(1.55), Inches(7.3), Inches(3.35), SURFACE)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(1.55), Inches(0.08), Inches(3.35))
    _fill_solid(bar, ACCENT)
    box = slide.shapes.add_textbox(Inches(0.65), Inches(1.7), Inches(6.9), Inches(3.05))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "📋 Passo a passo"
    _set_run_font(run, 14, True, ACCENT_HOVER)
    for i, passo in enumerate(a["passos"], start=1):
        _add_paragraph(tf, f"{i}. {passo}", size=13, color=TEXT, space_before=6, space_after=2)

    # Mock do WhatsApp
    _phone_mock(
        slide,
        Inches(8.0), Inches(1.55), Inches(4.9), Inches(3.35),
        "💬 Exemplo no WhatsApp",
        a["telas"],
    )

    # Dicas
    _add_round_rect(slide, Inches(0.4), Inches(5.05), Inches(6.1), Inches(1.9), TIP_BG)
    box = slide.shapes.add_textbox(Inches(0.55), Inches(5.15), Inches(5.8), Inches(1.7))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "💡 Dicas"
    _set_run_font(run, 12, True, ACCENT_HOVER)
    for d in a["dicas"]:
        _add_paragraph(tf, f"• {d}", size=11, color=TEXT, space_before=3)

    # Avisos
    _add_round_rect(slide, Inches(6.8), Inches(5.05), Inches(6.1), Inches(1.9), WARN_BG)
    box = slide.shapes.add_textbox(Inches(6.95), Inches(5.15), Inches(5.8), Inches(1.7))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "⚠️ Avisos importantes"
    _set_run_font(run, 12, True, WARNING)
    for av in a["avisos"]:
        _add_paragraph(tf, f"• {av}", size=11, color=TEXT, space_before=3)


def slide_vender_detalhe(prs: Presentation) -> None:
    """Slide extra só para VENDER — fluxo longo."""
    slide = _base_slide(prs)
    _header(slide, "VENDER — mapa do funil", "Do login PAP até a O.S. no CRM")

    etapas = [
        ("1", "Matrícula", "Confirma SIM"),
        ("2", "Endereço", "CEP · nº · ref."),
        ("3", "Cliente", "CPF · tel · e-mail"),
        ("4", "Crédito", "Análise PAP"),
        ("5", "Oferta", "Pgto · plano"),
        ("6", "Confirma", "SIM do cliente"),
        ("7", "Bio", "CONSULTAR/OK"),
        ("8", "Agenda", "Dia · período"),
        ("9", "O.S.", "Abre + CRM"),
    ]
    x0 = Inches(0.35)
    for i, (n, t, s) in enumerate(etapas):
        left = x0 + Inches(i * 1.42)
        _add_round_rect(slide, left, Inches(1.3), Inches(1.3), Inches(1.55), SURFACE)
        _add_rect(slide, left, Inches(1.3), Inches(1.3), Inches(0.35), ACCENT if i % 2 == 0 else SECONDARY)
        _textbox(slide, left, Inches(1.35), Inches(1.3), Inches(0.28), n, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _textbox(slide, left + Inches(0.05), Inches(1.75), Inches(1.2), Inches(0.45), t, size=11, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
        _textbox(slide, left + Inches(0.05), Inches(2.25), Inches(1.2), Inches(0.5), s, size=10, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
        if i < len(etapas) - 1:
            _textbox(slide, left + Inches(1.15), Inches(1.85), Inches(0.35), Inches(0.3), "→", size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    _card(
        slide, Inches(0.4), Inches(3.2), Inches(6.1), Inches(3.5),
        "Comandos auxiliares durante a venda",
        [
            "• CANCELAR / SAIR / PARAR — encerra o fluxo",
            "• ESTENDER — +5 minutos (máximo 3 vezes)",
            "• REPETIR — reenvia a última pergunta do bot",
            "• CONFIRMAR — quando o fluxo pedir confirmação",
            "• Timeout ≈ 10 min sem interação",
            "• Cliente recebe mensagem com protocolo",
        ],
        accent_bar=ACCENT,
    )
    _card(
        slide, Inches(6.8), Inches(3.2), Inches(6.1), Inches(3.5),
        "Checklist antes de iniciar",
        [
            "☑ Matrícula PAP válida e ativa",
            "☑ Flag autorizar_venda_sem_auditoria",
            "☑ Pool BackOffice PAP operacional",
            "☑ Cliente com WhatsApp no celular informado",
            "☑ CEP coberto e viável",
            "☑ Não misturar com NOVA VENDA (CRM)",
        ],
        title_color=SUCCESS,
        accent_bar=SUCCESS,
    )


def slide_permissoes(prs: Presentation) -> None:
    slide = _base_slide(prs)
    _header(slide, "Matriz rápida de permissões", "Quem pode usar o quê")

    rows = [
        ("Comando", "Quem usa", "Observação"),
        ("DFV, CDOE, Viabilidade", "Todos cadastrados", "Power BI / KMZ"),
        ("Status, Fatura, Conta", "Todos cadastrados", "CRM / Nio"),
        ("Material / Apoia", "Todos cadastrados", "Record Apoia"),
        ("Andamento", "Vend. (próprias) · Dir/Admin/BO (todas)", "Agenda do dia"),
        ("Inclusão", "Flag inclusão WPP", "Formulário externo"),
        ("Crédito / Pedido", "Flag análise crédito", "Pool BO PAP"),
        ("Bio", "Flag biometria", "Br Pronto / GED"),
        ("Vender", "Flag venda s/ auditoria + matrícula", "Funil PAP"),
        ("Nova Venda", "Cadastrados (+ flag automática)", "Só CRM"),
        ("Comissão", "Diretoria / Admin", "Lançamentos sensíveis"),
    ]

    top = Inches(1.2)
    col_w = [Inches(3.8), Inches(4.6), Inches(4.2)]
    headers_left = [Inches(0.4), Inches(4.2), Inches(8.8)]

    # header row
    for i, h in enumerate(rows[0]):
        _add_rect(slide, headers_left[i], top, col_w[i], Inches(0.38), ACCENT_HOVER)
        _textbox(slide, headers_left[i] + Inches(0.1), top + Inches(0.05), col_w[i] - Inches(0.15), Inches(0.3),
                 h, size=12, bold=True, color=WHITE)

    for ri, row in enumerate(rows[1:]):
        y = top + Inches(0.38) + Inches(ri * 0.48)
        bg = SURFACE if ri % 2 == 0 else BG_SECONDARY
        for i, cell in enumerate(row):
            _add_rect(slide, headers_left[i], y, col_w[i], Inches(0.48), bg)
            _textbox(slide, headers_left[i] + Inches(0.1), y + Inches(0.1), col_w[i] - Inches(0.15), Inches(0.35),
                     cell, size=11, color=TEXT)


def slide_atalhos(prs: Presentation) -> None:
    slide = _base_slide(prs)
    _header(slide, "Cola rápida — comandos", "Para fixar no mural ou no celular")

    cmds = [
        ("MENU", "Abre o menu completo"),
        ("DFV", "Fachadas por CEP"),
        ("CDOE", "Endereços por CDO"),
        ("VIABILIDADE", "Mancha / KMZ"),
        ("INCLUSAO", "Solicitar inclusão"),
        ("STATUS", "Status pedido/OS"),
        ("FATURA", "Faturas Nio"),
        ("CONTA", "2ª via site"),
        ("MATERIAL", "Buscar arquivo"),
        ("APOIA", "Record Apoia"),
        ("ANDAMENTO", "Agenda do dia"),
        ("CREDITO", "Análise crédito"),
        ("PEDIDO", "OS no PAP 30d"),
        ("BIO", "Biometria"),
        ("VENDER", "Venda no PAP"),
        ("NOVA VENDA", "Cadastro CRM"),
        ("COMISSAO", "Ajuda comissão"),
        ("CANCELAR", "Abortar fluxo"),
    ]

    for i, (cmd, desc) in enumerate(cmds):
        col = i % 3
        row = i // 3
        left = Inches(0.4 + col * 4.25)
        top = Inches(1.15 + row * 0.9)
        _add_round_rect(slide, left, top, Inches(4.0), Inches(0.75), SURFACE)
        _add_rect(slide, left, top, Inches(0.1), Inches(0.75), ACCENT if i % 2 == 0 else SECONDARY)
        _textbox(slide, left + Inches(0.25), top + Inches(0.08), Inches(3.6), Inches(0.3),
                 cmd, size=14, bold=True, color=ACCENT_HOVER)
        _textbox(slide, left + Inches(0.25), top + Inches(0.38), Inches(3.6), Inches(0.3),
                 desc, size=11, color=TEXT_MUTED)


def slide_encerramento(prs: Presentation) -> None:
    slide = _base_slide(prs, with_footer=False)
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, FOOTER)
    _add_rect(slide, 0, 0, Inches(0.25), SLIDE_H, ACCENT)

    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(0.7), Inches(1.5), height=Inches(0.7))

    _textbox(slide, Inches(0.7), Inches(2.6), Inches(11), Inches(0.6),
             "Pronto para operar no WhatsApp", size=32, bold=True, color=WHITE)
    _textbox(slide, Inches(0.7), Inches(3.4), Inches(11), Inches(1.2),
             "Em dúvida: digite MENU.\n"
             "Problema de permissão: fale com Governança.\n"
             "Instabilidade PAP/Nio: acione o BackOffice.\n\n"
             "Record PAP — parceiro oficial Nio Fibra.",
             size=16, color=RGBColor(0xCB, 0xD5, 0xE1))
    _add_rect(slide, Inches(0.7), Inches(5.5), Inches(3.5), Inches(0.08), ACCENT)
    _textbox(slide, Inches(0.7), Inches(5.8), Inches(10), Inches(0.4),
             "Obrigado — bom atendimento e boas vendas!",
             size=18, bold=True, color=ACCENT_LIGHT)


def main() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_capa(prs)
    slide_agenda(prs)
    slide_pre_requisitos(prs)

    # MATERIAL e APOIA já estão unidos no conteúdo
    seen_material = False
    idx = 0
    for a in AUTOMACOES:
        if a["cmd"].startswith("MATERIAL"):
            if seen_material:
                continue
            seen_material = True
        idx += 1
        slide_automacao(prs, idx, a)
        if a["cmd"] == "VENDER":
            slide_vender_detalhe(prs)

    slide_permissoes(prs)
    slide_atalhos(prs)
    slide_encerramento(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"OK: {OUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
