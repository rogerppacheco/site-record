# crm_app/conhecimento_ia_extract.py
"""
Extrai texto de PDF, Excel e PowerPoint para alimentar a base de conhecimento da IA.
"""
import logging
import os

logger = logging.getLogger(__name__)


def extrair_pdf(caminho: str) -> str:
    """Extrai texto de PDF (PyMuPDF)."""
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(caminho)
        partes = []
        for page in doc:
            partes.append(page.get_text())
        doc.close()
        return "\n".join(partes).strip()
    except Exception as e:
        logger.warning("[Conhecimento IA] Erro ao extrair PDF %s: %s", caminho, e)
        return ""


def extrair_excel(caminho: str) -> str:
    """Extrai texto de planilhas Excel (todas as abas e células)."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(caminho, read_only=True, data_only=True)
        partes = []
        for sheet in wb.worksheets:
            partes.append(f"[Aba: {sheet.title}]")
            for row in sheet.iter_rows(values_only=True):
                linha = "\t".join(str(c) if c is not None else "" for c in row)
                if linha.strip():
                    partes.append(linha)
        wb.close()
        return "\n".join(partes).strip()
    except Exception as e:
        logger.warning("[Conhecimento IA] Erro ao extrair Excel %s: %s", caminho, e)
        return ""


def extrair_ppt(caminho: str) -> str:
    """Extrai texto de PowerPoint (slides)."""
    try:
        from pptx import Presentation
        from pptx.util import Inches
        prs = Presentation(caminho)
        partes = []
        for i, slide in enumerate(prs.slides, 1):
            partes.append(f"[Slide {i}]")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    partes.append(shape.text.strip())
        return "\n".join(partes).strip()
    except ImportError:
        logger.warning("[Conhecimento IA] python-pptx não instalado. Instale com: pip install python-pptx")
        return ""
    except Exception as e:
        logger.warning("[Conhecimento IA] Erro ao extrair PPT %s: %s", caminho, e)
        return ""


def extrair_texto_arquivo(arquivo) -> str:
    """
    Recebe um objeto File (Django UploadedFile ou path) e extrai o texto.
    Retorna string vazia se tipo não suportado ou erro.
    """
    if isinstance(arquivo, str) and os.path.isfile(arquivo):
        return _extrair_por_extensao(arquivo, os.path.basename(arquivo))

    import tempfile
    nome = getattr(arquivo, "name", None) or os.path.basename(str(arquivo))
    ext = os.path.splitext(nome)[-1].lower()
    suf = ext if ext in (".pdf", ".xls", ".xlsx", ".ppt", ".pptx") else ".bin"

    if hasattr(arquivo, "chunks"):
        chunks = list(arquivo.chunks())
    elif hasattr(arquivo, "read"):
        chunks = [arquivo.read()]
    else:
        return ""

    with tempfile.NamedTemporaryFile(delete=False, suffix=suf) as tmp:
        for chunk in chunks:
            tmp.write(chunk)
        tmp.flush()
        caminho = tmp.name
    try:
        return _extrair_por_extensao(caminho, nome)
    finally:
        try:
            os.unlink(caminho)
        except OSError:
            pass


def _extrair_por_extensao(caminho: str, nome_arquivo: str) -> str:
    ext = os.path.splitext(nome_arquivo)[-1].lower()
    if ext == ".pdf":
        return extrair_pdf(caminho)
    if ext in (".xls", ".xlsx"):
        return extrair_excel(caminho)
    if ext in (".ppt", ".pptx"):
        return extrair_ppt(caminho)
    logger.warning("[Conhecimento IA] Tipo não suportado: %s", nome_arquivo)
    return ""
